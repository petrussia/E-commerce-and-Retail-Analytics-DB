import pptx
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


def making_presentation(total_sales_amount, total_profit, total_orders, new_customers, top_products_df, total_cancellations,
cancellation_impact, top_stores_by_profit, top_stores_by_sales ):
    # Создание презентации
    prs = pptx.Presentation()

    # Слайд 1: Название презентации
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Retail Analytical Metrics"
    subtitle.text = "Подготовлено для встречи с новым владельцем"

    # Слайд 2: Анализ клиентов
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Анализ клиентов"
    content = f"Общее количество уникальных клиентов у компании: {new_customers}"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = content

    # Слайд 3: Самые продаваемые товары
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Самые продаваемые товары"
    content = 'Топ-10 самых продаваемых товаров:\n'
    content += top_products_df.to_string(index=False)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4))
    tf = txBox.text_frame
    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line

    # Слайд 4: Анализ отменённых и возвращённых товаров
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Анализ отменённых и возвращённых товаров"
    content = f"Количество отменённых и возвращённых заказов: {total_cancellations}\n"
    content += f"Упущенная прибыль отменённых и возвращённых заказов: {cancellation_impact}"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = content

    # Слайд 5: Динамика продаж по дням
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Динамика продаж по дням'
    imgpath = 'daily_sales_line.png'
    slide.shapes.add_picture(imgpath, Inches(1), Inches(1.5), height=Inches(5))

    # Слайд 6: Динамика продаж по месяцам
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Динамика продаж по месяцам'
    imgpath = 'monthly_sales_pie.png'
    slide.shapes.add_picture(imgpath, Inches(1), Inches(1.5), height=Inches(5))

    # Слайд 7: Динамика продаж по дням недели, прибыли по месяцам
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Динамика продаж по дням недели и прибыли по месяцам'
    imgpath = 'weekly_sales_bar.png'
    slide.shapes.add_picture(imgpath, Inches(1), Inches(2), height=Inches(5))

    # Слайд 8: Динамика прибыли по месяцам
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Динамика продаж по дням недели и прибыли по месяцам'
    imgpath = 'monthly_profit_line.png'
    slide.shapes.add_picture(imgpath, Inches(1), Inches(2), height=Inches(5))

    # Слайд 9: Анализ магазинов по выручке
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Анализ магазинов по выручке"
    content = 'Топ-10 магазинов с наибольшей выручкой:\n'
    content += top_stores_by_sales.to_string()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4.5))
    tf = txBox.text_frame
    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line

    # Слайд 10: Анализ магазинов по прибыли
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Анализ магазинов по прибыли"
    content = 'Топ-10 магазинов с наибольшей прибылью:\n'
    content += top_stores_by_profit.to_string()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4.5))
    tf = txBox.text_frame
    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line

    # Слайд 11: Вывод общих метрик по всей сети
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Вывод общих метрик по всей сети"
    content = f"Общая выручка компании: {total_sales_amount}\n"
    content += f"Общая сумма прибыли компании: {total_profit}\n"
    content += f"Количество заказов у компании за всё время: {total_orders}\n"
    content += f"Количество новых клиентов: {new_customers}"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(2.5))
    tf = txBox.text_frame
    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line

    # Слайд 12: Ссылка на отчет с созданием гиперссылки на этот отчёт
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Ссылка на отчет"

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Отчет сохранен как 'Retails_analytics_of_company_report.csv'"
    r_id = slide.part.relate_to('Retails_analytics_of_company_report.csv', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hyperlink_field = OxmlElement('a:hlinkClick')
    hyperlink_field.set(qn('r:id'), r_id)
    run._r.append(hyperlink_field)

    # Сохранение презентации
    prs.save('Retail_Analytical_Metrics.pptx')
